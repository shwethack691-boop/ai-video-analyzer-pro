import os
import uuid
import whisper
import yt_dlp

from youtube_transcript_api import YouTubeTranscriptApi
from transformers import pipeline
from gtts import gTTS
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from docx import Document
from pptx import Presentation
from deep_translator import GoogleTranslator

# ---------------- MODELS ----------------
whisper_model = whisper.load_model("base")

summarizer = pipeline("summarization", model="facebook/bart-large-cnn")

qa_pipeline = pipeline("question-answering", model="deepset/roberta-base-squad2")

# ---------------- YOUTUBE TEXT ----------------
def get_youtube_text(url):
    try:
        video_id = url.split("v=")[1].split("&")[0]
        transcript = YouTubeTranscriptApi.get_transcript(video_id)

        text = " ".join([t["text"] for t in transcript])
        return text, transcript
    except:
        return None, None

# ---------------- DOWNLOAD AUDIO (FIXED) ----------------
def download_audio(url):
    file = f"audio_{uuid.uuid4()}.mp3"

    ydl_opts = {
        'format': 'bestaudio/best',
        'outtmpl': file,
        'quiet': True,
        'noplaylist': True
    }

    try:
        with yt_dlp.YoutubeDL(ydl_opts) as ydl:
            ydl.download([url])
        return file
    except Exception as e:
        raise Exception(f"Audio download failed: {e}")

# ---------------- TRANSCRIBE ----------------
def transcribe_audio(path):
    result = whisper_model.transcribe(path)
    text = result["text"]

    timestamps = []
    for seg in result.get("segments", []):
        timestamps.append({
            "start": seg["start"],
            "text": seg["text"]
        })

    return text, timestamps

# ---------------- SUMMARY ----------------
def summarize_text(text, mode="Medium"):
    length = {"Short": 100, "Medium": 200, "Long": 400}

    summary = summarizer(
        text[:2000],
        max_length=length[mode],
        min_length=50,
        do_sample=False
    )

    return summary[0]["summary_text"]

# ---------------- TRANSLATE ----------------
def translate_text(text, lang):
    try:
        return GoogleTranslator(source='auto', target=lang).translate(text[:3000])
    except:
        return text

# ---------------- HIGHLIGHTS ----------------
def highlight_text(text):
    sentences = text.split(".")
    return [s.strip() for s in sentences if len(s) > 60][:10]

# ---------------- AUTO CHAPTERS ----------------
def generate_chapters(timestamps):
    chapters = []
    current = ""
    start_time = 0

    for t in timestamps:
        if len(current) < 150:
            current += " " + t["text"]
        else:
            chapters.append({
                "time": round(start_time, 2),
                "title": current[:60]
            })
            current = t["text"]
            start_time = t["start"]

    if current:
        chapters.append({
            "time": round(start_time, 2),
            "title": current[:60]
        })

    return chapters

# ---------------- Q&A CHAT ----------------
def ask_question(context, question):
    try:
        answer = qa_pipeline({
            "question": question,
            "context": context[:3000]
        })
        return answer["answer"]
    except:
        return "Could not find answer."

# ---------------- AUDIO ----------------
def text_to_audio(text, lang="en"):
    file = f"audio_{uuid.uuid4()}.mp3"
    gTTS(text=text, lang=lang).save(file)
    return file

# ---------------- EXPORTS ----------------
def create_pdf(text):
    file = "output.pdf"
    doc = SimpleDocTemplate(file)
    styles = getSampleStyleSheet()
    content = [Paragraph(line, styles["Normal"]) for line in text.split("\n")]
    doc.build(content)
    return file

def create_docx(text):
    file = "output.docx"
    doc = Document()
    doc.add_heading("AI Video Report", 0)
    for line in text.split("\n"):
        doc.add_paragraph(line)
    doc.save(file)
    return file

def create_ppt(text):
    file = "output.pptx"
    prs = Presentation()
    for part in text.split("\n"):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "AI Slide"
        slide.placeholders[1].text = part[:200]
    prs.save(file)
    return file
